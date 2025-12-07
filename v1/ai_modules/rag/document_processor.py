"""
Procesador de Documentos (PDF y PowerPoint)
Piloto IA - FICEM BD

Extrae texto de PDFs y PowerPoint y los divide en chunks para embeddings.
"""

from pathlib import Path
from typing import List, Dict
import PyPDF2
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter


class DocumentProcessor:
    """
    Procesa documentos PDF y PowerPoint y los prepara para RAG.
    """

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Inicializa el procesador de documentos.

        Args:
            chunk_size: Tamaño de cada chunk en caracteres
            chunk_overlap: Solapamiento entre chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )

    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """
        Extrae todo el texto de un PDF.

        Args:
            pdf_path: Ruta al archivo PDF

        Returns:
            Texto completo del PDF
        """
        pdf_path = Path(pdf_path)

        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF no encontrado: {pdf_path}")

        text = ""

        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)

                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"\n--- Página {page_num + 1} ---\n{page_text}"

        except Exception as e:
            print(f"Error leyendo PDF {pdf_path.name}: {e}")
            return ""

        return text

    def extract_text_from_pptx(self, pptx_path: str) -> str:
        """
        Extrae todo el texto de un PowerPoint.

        Args:
            pptx_path: Ruta al archivo PPTX

        Returns:
            Texto completo del PowerPoint
        """
        pptx_path = Path(pptx_path)

        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint no encontrado: {pptx_path}")

        text = ""

        try:
            prs = Presentation(str(pptx_path))

            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Diapositiva {slide_num + 1} ---\n"

                # Extraer texto de todas las formas
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"

                # Extraer texto de notas
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text:
                        slide_text += f"\n[Notas: {notes_frame.text}]\n"

                text += slide_text

        except Exception as e:
            print(f"Error leyendo PowerPoint {pptx_path.name}: {e}")
            return ""

        return text

    def process_pdf(self, pdf_path: str, metadata: Dict = None) -> List[Dict]:
        """
        Procesa un PDF completo: extrae texto y crea chunks.

        Args:
            pdf_path: Ruta al archivo PDF
            metadata: Metadatos adicionales para agregar a cada chunk

        Returns:
            Lista de chunks con texto y metadatos
        """
        pdf_path = Path(pdf_path)

        # Extraer texto
        text = self.extract_text_from_pdf(str(pdf_path))

        if not text:
            print(f"⚠️  No se pudo extraer texto de {pdf_path.name}")
            return []

        # Dividir en chunks
        chunks = self.text_splitter.split_text(text)

        # Agregar metadatos
        base_metadata = {
            "source": pdf_path.name,
            "source_path": str(pdf_path),
            "total_chunks": len(chunks)
        }

        if metadata:
            base_metadata.update(metadata)

        # Crear documentos con metadatos
        documents = []
        for i, chunk in enumerate(chunks):
            doc_metadata = base_metadata.copy()
            doc_metadata["chunk_id"] = i

            documents.append({
                "text": chunk,
                "metadata": doc_metadata
            })

        print(f"✅ Procesado {pdf_path.name}: {len(chunks)} chunks")
        return documents

    def process_pptx(self, pptx_path: str, metadata: Dict = None) -> List[Dict]:
        """
        Procesa un PowerPoint completo: extrae texto y crea chunks.

        Args:
            pptx_path: Ruta al archivo PPTX
            metadata: Metadatos adicionales para agregar a cada chunk

        Returns:
            Lista de chunks con texto y metadatos
        """
        pptx_path = Path(pptx_path)

        # Extraer texto
        text = self.extract_text_from_pptx(str(pptx_path))

        if not text:
            print(f"⚠️  No se pudo extraer texto de {pptx_path.name}")
            return []

        # Dividir en chunks
        chunks = self.text_splitter.split_text(text)

        # Agregar metadatos
        base_metadata = {
            "source": pptx_path.name,
            "source_path": str(pptx_path),
            "total_chunks": len(chunks)
        }

        if metadata:
            base_metadata.update(metadata)

        # Crear documentos con metadatos
        documents = []
        for i, chunk in enumerate(chunks):
            doc_metadata = base_metadata.copy()
            doc_metadata["chunk_id"] = i

            documents.append({
                "text": chunk,
                "metadata": doc_metadata
            })

        print(f"✅ Procesado {pptx_path.name}: {len(chunks)} chunks")
        return documents

    def process_directory(self, directory_path: str, file_pattern: str = "*.pdf") -> List[Dict]:
        """
        Procesa todos los archivos (PDF o PPTX) en un directorio.

        Args:
            directory_path: Ruta al directorio
            file_pattern: Patrón de archivos a procesar (*.pdf, *.pptx, o *.*)

        Returns:
            Lista de todos los chunks de todos los documentos
        """
        directory = Path(directory_path)

        if not directory.exists():
            raise FileNotFoundError(f"Directorio no encontrado: {directory}")

        all_documents = []
        files = list(directory.glob(file_pattern))

        print(f"📂 Procesando {len(files)} archivo(s)...")

        for file in files:
            # Metadatos específicos por documento
            metadata = {
                "document_type": self._classify_document(file.name),
                "file_size_mb": file.stat().st_size / (1024 * 1024)
            }

            # Procesar según extensión
            if file.suffix.lower() == '.pdf':
                documents = self.process_pdf(str(file), metadata)
            elif file.suffix.lower() == '.pptx':
                documents = self.process_pptx(str(file), metadata)
            else:
                print(f"⚠️  Tipo de archivo no soportado: {file.name}")
                continue

            all_documents.extend(documents)

        print(f"✅ Total: {len(all_documents)} chunks de {len(files)} documentos")
        return all_documents

    def _classify_document(self, filename: str) -> str:
        """
        Clasifica el tipo de documento según el nombre.

        Args:
            filename: Nombre del archivo

        Returns:
            Tipo de documento
        """
        filename_lower = filename.lower()

        if "gnr" in filename_lower and "concrete" in filename_lower:
            return "gcca_gnr_concrete"
        elif "gcca" in filename_lower:
            return "gcca_roadmap"
        elif "ecra" in filename_lower:
            return "ecra_technology"
        elif "cembureau" in filename_lower:
            return "cembureau_roadmap"
        elif "roadmap" in filename_lower:
            return "industry_roadmap"
        elif "peru" in filename_lower or "perú" in filename_lower:
            return "peru_roadmap"
        elif "mpa" in filename_lower or "uk" in filename_lower:
            return "uk_roadmap"
        else:
            return "technical_document"


# Ejemplo de uso
if __name__ == "__main__":
    processor = DocumentProcessor(chunk_size=1000, chunk_overlap=200)

    # Procesar un PDF
    docs = processor.process_directory("docs/tech_docs")

    print(f"\nEjemplo de chunk:")
    print(docs[0])
